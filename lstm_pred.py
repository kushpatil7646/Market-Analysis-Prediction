import os, json, math
from datetime import datetime
# ----------------- CONFIG -----------------
FORCE_CPU = True             # True to force CPU (safe on mismatched CUDA/cuDNN)
PLOT_DPI = 300               # high resolution for plots
RANDOM_SEED = 42             # reproducibility
MODEL_PATH = "best_model.keras"
FORECAST_HORIZON = 6         # months ahead
EXOG_WINDOW = 24             # months used to extrapolate price/export
EXOG_METHOD = "linreg"       # "linreg" or "hold"
OUTPUT_DIR = "outputs"       # where ALL artifacts go
REPORT_NAME = "Executive_Summary.pptx"
# ------------------------------------------

if FORCE_CPU:
    os.environ["CUDA_VISIBLE_DEVICES"] = "-1"

import numpy as np
import pandas as pd
from sklearn.preprocessing import MinMaxScaler
from sklearn.linear_model import LinearRegression
import tensorflow as tf
from tensorflow.keras import Input, Sequential
from tensorflow.keras.layers import LSTM, Dense, Dropout
from tensorflow.keras.callbacks import EarlyStopping, ModelCheckpoint
import matplotlib.pyplot as plt
from sklearn.metrics import mean_absolute_error, mean_squared_error
from pandas.tseries.offsets import MonthEnd
import joblib

# Optional report packaging
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    HAVE_PPTX = True
except Exception:
    HAVE_PPTX = False

# Reproducibility
np.random.seed(RANDOM_SEED)
tf.random.set_seed(RANDOM_SEED)

# ---------- IO Helpers ----------
os.makedirs(OUTPUT_DIR, exist_ok=True)
STAMP = datetime.now().strftime("%Y%m%d_%H%M%S")

def out_path(name):
    base, ext = os.path.splitext(name)
    return os.path.join(OUTPUT_DIR, f"{base}_{STAMP}{ext}")

def save_fig(fig, filename):
    path = out_path(filename)
    fig.savefig(path, format=filename.split(".")[-1], dpi=PLOT_DPI, bbox_inches="tight")
    plt.close(fig)
    print(f"Saved: {path}")
    return path

# ---------- Generate More Realistic Dummy Sales Data ----------
def generate_dummy_data():
    periods = 1000
    dates = pd.date_range(start="2018-01-01", periods=periods, freq="M")
    time = np.arange(periods)
    trend = 0.15 * time
    seasonality = 25 * np.sin(2 * np.pi * time / 12)  # yearly seasonality
    noise = np.random.normal(0, 8, periods)
    demand = np.abs(200 + trend + seasonality + noise)
    price = 50 + 0.2 * demand + np.random.normal(0, 1.5, periods)
    export_index = 100 + 0.5 * time + 10 * np.sin(2 * np.pi * time / 24) + np.random.normal(0, 3, periods)
    return pd.DataFrame({
        "date": dates,
        "demand": np.round(demand, 2),
        "price": np.round(price, 2),
        "export_index": np.round(export_index, 2)
    })

CSV_PATH = "sales_data.csv"
if not os.path.exists(CSV_PATH):
    df_dummy = generate_dummy_data()
    df_dummy.to_csv(CSV_PATH, index=False)
    print(f"Dummy sales data created: {CSV_PATH}")

# ---------- Load Data ----------
df = pd.read_csv(CSV_PATH, parse_dates=["date"])
df = df.sort_values("date").reset_index(drop=True)

# ---------- Feature Engineering ----------
features = ["demand", "price", "export_index"]
scaler = MinMaxScaler()
scaled = scaler.fit_transform(df[features])

LOOKBACK = 12
def create_sequences(X, lookback):
    Xs, ys = [], []
    for i in range(len(X) - lookback):
        Xs.append(X[i:i+lookback])
        ys.append(X[i+lookback, 0])  # demand is first col
    return np.array(Xs), np.array(ys)

X, y = create_sequences(scaled, LOOKBACK)

# Train/val/test split
test_size = int(len(X) * 0.2)
val_size = int(len(X) * 0.1)

X_train, X_test = X[:-test_size], X[-test_size:]
y_train, y_test = y[:-test_size], y[-test_size:]

X_train, X_val = X_train[:-val_size], X_train[-val_size:]
y_train, y_val = y_train[:-val_size], y_train[-val_size:]

# Memory optimization
X_train = np.ascontiguousarray(X_train.astype("float32"))
X_val   = np.ascontiguousarray(X_val.astype("float32"))
X_test  = np.ascontiguousarray(X_test.astype("float32"))
y_train = np.ascontiguousarray(y_train.astype("float32").reshape(-1, 1))
y_val   = np.ascontiguousarray(y_val.astype("float32").reshape(-1, 1))
y_test  = np.ascontiguousarray(y_test.astype("float32").reshape(-1, 1))

# ---------- Model ----------
model = Sequential([
    Input(shape=(X_train.shape[1], X_train.shape[2])),
    LSTM(64, return_sequences=True),
    Dropout(0.2),
    LSTM(32),
    Dropout(0.2),
    Dense(1)
])

model.compile(optimizer=tf.keras.optimizers.Adam(learning_rate=0.001), 
              loss="mse",
              metrics=["mae"])

# Callbacks
es = EarlyStopping(monitor="val_loss", patience=50, restore_best_weights=True, verbose=1)
mc = ModelCheckpoint(MODEL_PATH, monitor="val_loss", save_best_only=True)

history = model.fit(
    X_train, y_train,
    validation_data=(X_val, y_val),
    epochs=50,
    batch_size=16,
    callbacks=[es, mc],
    verbose=1
)

# Load the best model
best_model = tf.keras.models.load_model(MODEL_PATH)

# ---------- Evaluation ----------
def invert_target(scaled_target):
    # inverse only for demand column (0)
    dummy = np.zeros((len(scaled_target), scaled.shape[1]))
    dummy[:, 0] = scaled_target.flatten()
    return scaler.inverse_transform(dummy)[:, 0]

y_pred = best_model.predict(X_test, verbose=0)
y_test_inv = invert_target(y_test)
y_pred_inv = invert_target(y_pred)

mae = mean_absolute_error(y_test_inv, y_pred_inv)
rmse = math.sqrt(mean_squared_error(y_test_inv, y_pred_inv))
mape = np.mean(np.abs((y_test_inv - y_pred_inv) / np.maximum(y_test_inv, 1e-6))) * 100.0

print(f"MAE: {mae:.2f} | RMSE: {rmse:.2f} | MAPE: {mape:.2f}%")

# ---------- Helper: Exogenous Projection ----------
def project_exogenous(df, horizon=6, method="linreg", window=24):
    """
    Forecast price and export_index for `horizon` months.
    method: "linreg" (linear trend on last `window` points) or "hold" (last value).
    Returns arrays: future_dates, price_future, export_future
    """
    last_date = df["date"].iloc[-1]
    future_dates = pd.date_range(last_date + MonthEnd(1), periods=horizon, freq="M")

    def linreg_extrap(series, steps, win):
        y = series.values[-win:].reshape(-1, 1)
        x = np.arange(len(y)).reshape(-1, 1)
        model = LinearRegression().fit(x, y)
        x_future = np.arange(len(y), len(y)+steps).reshape(-1, 1)
        y_future = model.predict(x_future).flatten()
        return y_future

    if method == "linreg":
        price_future = linreg_extrap(df["price"], horizon, min(window, len(df)))
        export_future = linreg_extrap(df["export_index"], horizon, min(window, len(df)))
    else:  # "hold"
        price_future = np.repeat(df["price"].iloc[-1], horizon)
        export_future = np.repeat(df["export_index"].iloc[-1], horizon)

    return future_dates, price_future, export_future

# ---------- 6-Month Forecast (Recursive, multivariate) ----------
def scale_value(col_index, value):
    # MinMaxScaler: X_scaled = X * scale_ + min_
    return value * scaler.scale_[col_index] + scaler.min_[col_index]

def six_month_forecast(best_model, df, scaled, horizon=6, method="linreg", window=24):
    # Prepare future exogenous (raw)
    future_dates, price_f_raw, export_f_raw = project_exogenous(df, horizon, method, window)

    # Scale exogenous using the scaler per-column formula
    price_f_scaled = np.array([scale_value(1, v) for v in price_f_raw])
    export_f_scaled = np.array([scale_value(2, v) for v in export_f_raw])

    # Start with the last LOOKBACK window
    window_seq = scaled[-LOOKBACK:].copy()  # shape (LOOKBACK, 3)
    preds_scaled = []

    for i in range(horizon):
        # Predict next demand (scaled)
        next_demand_scaled = best_model.predict(window_seq[np.newaxis, ...], verbose=0)[0, 0]
        preds_scaled.append(next_demand_scaled)

        # Compose next timestep (scaled feature row)
        next_row = np.array([next_demand_scaled, price_f_scaled[i], export_f_scaled[i]], dtype=np.float32)

        # Slide the window
        window_seq = np.vstack([window_seq[1:], next_row])

    # Invert to original units
    preds_inv = invert_target(np.array(preds_scaled).reshape(-1, 1))
    return pd.DataFrame({
        "date": future_dates,
        "forecast_demand": np.round(preds_inv, 2),
        "projected_price": np.round(price_f_raw, 2),
        "projected_export_index": np.round(export_f_raw, 2)
    })

forecast_df = six_month_forecast(
    best_model, df, scaled,
    horizon=FORECAST_HORIZON,
    method=EXOG_METHOD,
    window=EXOG_WINDOW
)

# ---------- Save Data Artifacts ----------
forecast_csv = out_path("six_month_forecast.csv")
forecast_df.to_csv(forecast_csv, index=False)

history_df = pd.DataFrame({
    "epoch": np.arange(1, len(history.history["loss"])+1),
    "loss": history.history["loss"],
    "val_loss": history.history["val_loss"],
    "mae": history.history.get("mae", []),
    "val_mae": history.history.get("val_mae", [])
})
history_csv = out_path("training_history.csv")
history_df.to_csv(history_csv, index=False)

corr_df = df[features].corr()
corr_csv = out_path("feature_correlations.csv")
corr_df.to_csv(corr_csv)

kpis = {
    "timestamp": STAMP,
    "records": int(len(df)),
    "lookback_months": int(LOOKBACK),
    "forecast_horizon_months": int(FORECAST_HORIZON),
    "exog_method": EXOG_METHOD,
    "exog_window": int(EXOG_WINDOW),
    "metrics": {
        "MAE": round(float(mae), 3),
        "RMSE": round(float(rmse), 3),
        "MAPE_percent": round(float(mape), 3)
    }
}
kpi_json = out_path("kpis.json")
with open(kpi_json, "w") as f:
    json.dump(kpis, f, indent=2)
print(f"Saved data: {forecast_csv}, {history_csv}, {corr_csv}, {kpi_json}")

# Save the scaler for future use
scaler_path = out_path("scaler.save")
joblib.dump(scaler, scaler_path)
print(f"Scaler saved to {scaler_path}")

# ---------- Visualizations (Saved to disk) ----------
# Plot 1: Training History
fig1 = plt.figure(figsize=(12, 6))
plt.plot(history.history["loss"], label="Training Loss")
plt.plot(history.history["val_loss"], label="Validation Loss")
plt.title("Model Training History")
plt.xlabel("Epoch"); plt.ylabel("Loss (MSE)")
plt.legend(); plt.grid(True)
fig1_path = save_fig(fig1, "training_history.png")

# Plot 2: Actual vs Predicted (test window)
test_dates = df["date"][-len(y_test):]
fig2 = plt.figure(figsize=(14, 7))
plt.plot(test_dates, y_test_inv, label="Actual Demand", marker="o", markersize=4)
plt.plot(test_dates, y_pred_inv, label="Predicted Demand", linestyle="--", marker="x", markersize=5)
plt.title("Demand Forecast: Actual vs Predicted (Test Set)")
plt.xlabel("Date"); plt.ylabel("Demand")
plt.legend(); plt.grid(True); plt.xticks(rotation=45)
fig2_path = save_fig(fig2, "demand_forecast_test.png")

# Plot 3: Feature Correlations
fig3 = plt.figure(figsize=(8, 6))
plt.matshow(corr_df.values, fignum=1, cmap="coolwarm")
plt.colorbar()
plt.xticks(range(len(features)), features, rotation=45)
plt.yticks(range(len(features)), features)
plt.title("Feature Correlation Matrix")
fig3_path = save_fig(fig3, "feature_correlations.png")

# Plot 4: History + 6-Month Forecast (last 36 months + forecast)
fig4 = plt.figure(figsize=(14, 7))
hist_tail = df.tail(36)
plt.plot(hist_tail["date"], hist_tail["demand"], label="Historical Demand", marker="o", markersize=3)
plt.plot(forecast_df["date"], forecast_df["forecast_demand"], label="6-Month Forecast", linestyle="--", marker="x", markersize=5)
plt.title("Historical Demand with 6-Month Ahead Forecast")
plt.xlabel("Date"); plt.ylabel("Demand")
plt.legend(); plt.grid(True); plt.xticks(rotation=45)
fig4_path = save_fig(fig4, "history_plus_6mo_forecast.png")

# ---------- Optional Excel Pack (nice for corporate reviewers) ----------
try:
    import xlsxwriter
    excel_path = out_path("Executive_Data_Pack.xlsx")
    with pd.ExcelWriter(excel_path, engine="xlsxwriter") as writer:
        df.to_excel(writer, sheet_name="Raw_Data", index=False)
        forecast_df.to_excel(writer, sheet_name="6M_Forecast", index=False)
        corr_df.to_excel(writer, sheet_name="Correlations")
        pd.DataFrame([kpis["metrics"]]).to_excel(writer, sheet_name="KPIs", index=False)
    print(f"Saved Excel pack: {excel_path}")
except Exception as e:
    print(f"(Optional) Excel pack not created: {e}")

# ---------- Optional PPT Executive Summary ----------
def make_pptx(kpis, fig_paths, forecast_csv_path, ppt_name):
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Demand Forecast Executive Summary"
    slide.placeholders[1].text = f"Run: {kpis['timestamp']}\nRecords: {kpis['records']}\nLookback: {kpis['lookback_months']} | Horizon: {kpis['forecast_horizon_months']}"

    # KPI slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tf = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.4), Inches(1.2)).text_frame
    tf.text = "Key Performance Indicators"
    p = tf.add_paragraph()
    p.text = f"RMSE: {kpis['metrics']['RMSE']} | MAE: {kpis['metrics']['MAE']} | MAPE: {kpis['metrics']['MAPE_percent']}%"
    p.level = 1

    # Charts slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    left = Inches(0.5); top = Inches(1.0)
    width = Inches(4.6)
    # place up to 3 charts
    for i, path in enumerate(fig_paths[:3]):
        slide.shapes.add_picture(path, left + Inches((i%2)*4.8), top + Inches((i//2)*3.6), width=width)

    # Forecast chart slide (full width)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_picture(fig_paths[-1], Inches(0.5), Inches(1.0), width=Inches(9.0))
    # Footer
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9.0), Inches(0.5)).text_frame
    tx.text = f"Forecast details: {os.path.basename(forecast_csv_path)}"

    # Save
    outp = out_path(ppt_name)
    prs.save(outp)
    print(f"PPT saved: {outp}")
    return outp

ppt_path = None
if HAVE_PPTX:
    fig_paths = [fig1_path, fig2_path, fig3_path, fig4_path]
    ppt_path = make_pptx(kpis, fig_paths, forecast_csv, REPORT_NAME)
else:
    print("(Optional) PPTX not created: install `python-pptx` to enable.")

# ---------- Console Summary ----------
print("\n=== MODEL EVALUATION ===")
print(f"MAE: {mae:.2f}")
print(f"RMSE: {rmse:.2f}")
print(f"MAPE: {mape:.2f}%")
print("\nArtifacts saved in:", os.path.abspath(OUTPUT_DIR))
